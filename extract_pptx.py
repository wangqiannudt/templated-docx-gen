from pptx import Presentation
import sys
prs = Presentation(sys.argv[1])
with open(sys.argv[2], 'w', encoding='utf-8') as f:
    f.write(f"=== 总幻灯片数: {len(prs.slides)} ===\n\n")
    for i, slide in enumerate(prs.slides):
        f.write(f"\n========== 第 {i+1} 页 ==========\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        f.write(text + "\n")
            if shape.has_table:
                f.write("[表格]\n")
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace('\n',' ') for cell in row.cells]
                    f.write(" | ".join(cells) + "\n")
print(f"提取完成: {sys.argv[2]}, 共 {len(prs.slides)} 页")
